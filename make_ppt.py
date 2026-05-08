"""Generate clean DBMS Presentation — No code, only database concepts & diagrams"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(15, 23, 42)
CARD = RGBColor(30, 41, 69)
ACCENT = RGBColor(99, 145, 255)
CYAN = RGBColor(56, 189, 248)
WHITE = RGBColor(240, 240, 250)
MUTED = RGBColor(148, 163, 184)
GREEN = RGBColor(52, 211, 153)
ORANGE = RGBColor(251, 146, 60)
RED = RGBColor(251, 113, 133)
GOLD = RGBColor(250, 204, 21)

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def txt(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = clr; p.font.bold = bold; p.font.name = 'Calibri'; p.alignment = align
    return tf

def para(tf, text, sz=16, clr=WHITE, bold=False, align=PP_ALIGN.LEFT, spacing=Pt(4)):
    p = tf.add_paragraph(); p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = clr; p.font.bold = bold; p.alignment = align
    p.space_before = spacing
    return p

def card(slide, l, t, w, h, fill=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = RGBColor(55, 65, 100); s.line.width = Pt(1)
    return s

def bullet_card(slide, x, y, w, h, title, items, icon="", title_clr=ACCENT):
    card(slide, x, y, w, h)
    txt(slide, x + 0.2, y + 0.15, w - 0.4, 0.4, f"{icon}  {title}" if icon else title, sz=18, clr=title_clr, bold=True, align=PP_ALIGN.CENTER)
    tf = txt(slide, x + 0.3, y + 0.65, w - 0.6, h - 0.8, "", sz=14, clr=WHITE)
    for item in items:
        para(tf, f"▸  {item}", sz=14, clr=WHITE, spacing=Pt(6))

# ==================== SLIDE 1: TITLE ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 1.0, 13.333, 0.8, "🎓", sz=56, align=PP_ALIGN.CENTER)
txt(s, 0, 2.0, 13.333, 0.8, "Smart College Event Management System", sz=40, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0, 3.0, 13.333, 0.6, "Database Design & Methodology", sz=26, clr=ACCENT, align=PP_ALIGN.CENTER)
txt(s, 0, 4.0, 13.333, 0.5, "DBMS Mini Project", sz=20, clr=MUTED, align=PP_ALIGN.CENTER)
card(s, 3.5, 5.2, 6.3, 1.2)
tf = txt(s, 3.8, 5.3, 5.7, 1, "4 Tables  •  BCNF Normalized  •  ER Model  •  Referential Integrity", sz=15, clr=CYAN, align=PP_ALIGN.CENTER)
para(tf, "Triggers  •  Views  •  Transactions  •  Indexes  •  ACID Properties", sz=15, clr=ACCENT, align=PP_ALIGN.CENTER)

# ==================== SLIDE 2: PROBLEM STATEMENT ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.3, 13.333, 0.6, "📌  Problem Statement", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)

card(s, 1.5, 1.3, 10.3, 2.0)
tf = txt(s, 1.8, 1.4, 9.7, 1.8, "", sz=16, clr=WHITE)
para(tf, "Colleges conduct numerous events (workshops, fests, seminars) throughout the year.", sz=16, clr=WHITE)
para(tf, "Managing registrations manually leads to:", sz=16, clr=WHITE, spacing=Pt(12))
para(tf, "▸  Overbooking — more students registered than seats available", sz=15, clr=ORANGE)
para(tf, "▸  Data redundancy — same information stored in multiple places", sz=15, clr=ORANGE)
para(tf, "▸  No tracking — no record of who registered for what", sz=15, clr=ORANGE)

txt(s, 0, 3.7, 13.333, 0.5, "Our Solution", sz=22, clr=GREEN, bold=True, align=PP_ALIGN.CENTER)

bullet_card(s, 0.8, 4.3, 3.8, 2.8, "For Students", ["Register using USN (unique ID)", "Browse & search events", "Register/unregister for events", "View their registrations"], icon="👨‍🎓", title_clr=CYAN)
bullet_card(s, 4.8, 4.3, 3.8, 2.8, "For Admins", ["Create, edit, delete events", "View registrations per event", "Manage student records", "See participation summary"], icon="🛡️", title_clr=CYAN)
bullet_card(s, 8.8, 4.3, 3.8, 2.8, "Database Goals", ["Eliminate data redundancy", "Ensure data integrity", "Prevent overbooking automatically", "ACID-compliant transactions"], icon="🗄️", title_clr=CYAN)

# ==================== SLIDE 3: ER DIAGRAM ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.2, 13.333, 0.6, "📐  Entity-Relationship (ER) Diagram", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)
er = os.path.join(os.path.dirname(__file__), 'static', 'er_diagram.png')
if os.path.exists(er):
    s.shapes.add_picture(er, Inches(0.5), Inches(1.0), Inches(6.5), Inches(5.8))

bullet_card(s, 7.5, 1.0, 5.3, 2.6, "Entities", ["STUDENT — stores student data (USN, name, dept)", "FACULTY — stores admin data", "EVENT — stores event details (name, date, venue)", "REGISTRATION — links students to events"])
bullet_card(s, 7.5, 3.8, 5.3, 1.5, "Relationships", ["Faculty ORGANIZES Events → 1 : N", "Students REGISTER FOR Events → M : N"])
bullet_card(s, 7.5, 5.5, 5.3, 1.5, "Key Observation", ["M:N relationship resolved via REGISTRATION table", "Acts as a junction/bridge table"])

# ==================== SLIDE 4: SCHEMA DIAGRAM ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.2, 13.333, 0.6, "🗂️  Relational Schema", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)
sch = os.path.join(os.path.dirname(__file__), 'static', 'schema_diagram.png')
if os.path.exists(sch):
    s.shapes.add_picture(sch, Inches(0.5), Inches(1.0), Inches(6.5), Inches(5.8))

bullet_card(s, 7.5, 1.0, 5.3, 2.0, "Primary Keys (Blue)", ["Student_ID — uniquely identifies each student", "Faculty_ID — uniquely identifies each admin", "Event_ID — uniquely identifies each event", "Reg_ID — uniquely identifies each registration"], title_clr=CYAN)
bullet_card(s, 7.5, 3.2, 5.3, 2.2, "Foreign Keys (Orange)", ["EVENT.Faculty_ID → FACULTY (ON DELETE SET NULL)", "REG.Student_ID → STUDENT (ON DELETE CASCADE)", "REG.Event_ID → EVENT (ON DELETE CASCADE)", "Ensures referential integrity across tables"], title_clr=ORANGE)
bullet_card(s, 7.5, 5.6, 5.3, 1.4, "Unique Constraints", ["Student USN — primary identity for students", "Student & Faculty Email — prevents duplicates", "UNIQUE(Student_ID, Event_ID) — no double registration"], title_clr=GREEN)

# ==================== SLIDE 5: NORMALIZATION ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.3, 13.333, 0.6, "✅  Normalization — BCNF", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0, 0.9, 13.333, 0.4, "All tables satisfy Boyce-Codd Normal Form, eliminating all anomalies", sz=15, clr=MUTED, align=PP_ALIGN.CENTER)

norms = [
    ("1NF", GREEN, "First Normal Form", ["All attribute values are atomic", "No multi-valued attributes", "No repeating groups", "Each cell has a single value"]),
    ("2NF", CYAN, "Second Normal Form", ["Satisfies 1NF", "No partial dependencies", "Non-key attributes depend on", "  the ENTIRE primary key"]),
    ("3NF", ACCENT, "Third Normal Form", ["Satisfies 2NF", "No transitive dependencies", "Non-key attributes depend", "  ONLY on primary key"]),
    ("BCNF", GOLD, "Boyce-Codd NF", ["Satisfies 3NF", "Every determinant is a", "  candidate key", "Strictest practical form"]),
]
for i, (label, clr, full, items) in enumerate(norms):
    x = 0.6 + i * 3.15
    card(s, x, 1.5, 2.9, 4.0)
    txt(s, x, 1.6, 2.9, 0.5, label, sz=28, clr=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, 2.1, 2.9, 0.3, full, sz=12, clr=MUTED, align=PP_ALIGN.CENTER)
    tf = txt(s, x + 0.2, 2.6, 2.5, 2.5, "", sz=13, clr=WHITE)
    for item in items:
        para(tf, f"▸  {item}" if not item.startswith("  ") else f"   {item}", sz=13, clr=WHITE, spacing=Pt(5))

card(s, 1.5, 5.8, 10.3, 1.3)
tf = txt(s, 1.8, 5.9, 9.7, 0.5, "Why Normalize?", sz=16, clr=GOLD, bold=True)
para(tf, "Without normalization, databases suffer from anomalies:", sz=14, clr=MUTED)
para(tf, "▸  Insertion Anomaly — cannot add data without unrelated data", sz=13, clr=WHITE, spacing=Pt(4))
para(tf, "▸  Update Anomaly — inconsistency when same data stored in multiple rows", sz=13, clr=WHITE, spacing=Pt(2))
para(tf, "▸  Deletion Anomaly — deleting one record accidentally removes other needed data", sz=13, clr=WHITE, spacing=Pt(2))

# ==================== SLIDE 6: INTEGRITY & CONSTRAINTS ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.3, 13.333, 0.6, "🔒  Data Integrity & Constraints", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)

bullet_card(s, 0.6, 1.2, 4.0, 2.8, "Entity Integrity", ["Every table has a PRIMARY KEY", "Primary keys are auto-incremented", "PRIMARY KEY = NOT NULL + UNIQUE", "Guarantees every row is identifiable"], icon="🔑", title_clr=CYAN)
bullet_card(s, 4.8, 1.2, 4.0, 2.8, "Referential Integrity", ["Foreign keys reference valid primary keys", "ON DELETE CASCADE — auto-remove child rows", "ON DELETE SET NULL — preserve but unlink", "No orphan records possible"], icon="🔗", title_clr=ORANGE)
bullet_card(s, 8.9, 1.2, 4.0, 2.8, "Domain Integrity", ["NOT NULL — required fields enforced", "UNIQUE — no duplicate USN or email", "DEFAULT values — Max_Seats = 100", "TEXT, INTEGER types enforce format"], icon="🛡️", title_clr=GREEN)

card(s, 0.6, 4.3, 12.2, 2.8)
txt(s, 0.9, 4.4, 11.6, 0.4, "Constraints Applied in Our Schema", sz=18, clr=GOLD, bold=True)
constraints = [
    ("PRIMARY KEY", "Student_ID, Faculty_ID, Event_ID, Reg_ID", "Unique row identifier"),
    ("FOREIGN KEY", "EVENT → FACULTY, REGISTRATION → STUDENT, EVENT", "Referential links"),
    ("UNIQUE", "Student USN, Student Email, Faculty Email", "No duplicates"),
    ("NOT NULL", "Name, Email, Password, USN, Date, Venue", "Required fields"),
    ("COMPOSITE UNIQUE", "UNIQUE(Student_ID, Event_ID) in REGISTRATION", "Prevents double registration"),
    ("DEFAULT", "Max_Seats DEFAULT 100, Timestamp DEFAULT CURRENT_TIMESTAMP", "Auto-fill values"),
]
for j, (ctype, applied, purpose) in enumerate(constraints):
    y = 5.0 + j * 0.32
    tf = txt(s, 1.0, y, 11, 0.3, "", sz=12, clr=WHITE)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"  {ctype}"; r1.font.color.rgb = CYAN; r1.font.bold = True; r1.font.size = Pt(12)
    r2 = p.add_run(); r2.text = f"   →   {applied}"; r2.font.color.rgb = WHITE; r2.font.size = Pt(12)
    r3 = p.add_run(); r3.text = f"   ({purpose})"; r3.font.color.rgb = MUTED; r3.font.size = Pt(11)

# ==================== SLIDE 7: DB FEATURES — TRIGGER & VIEW ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.3, 13.333, 0.6, "⚡  Database Features — Trigger & View", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)

card(s, 0.6, 1.2, 6.0, 5.5)
txt(s, 0.9, 1.3, 5.4, 0.5, "⚡  Trigger: Prevent Overbooking", sz=20, clr=ORANGE, bold=True)
tf = txt(s, 0.9, 1.9, 5.4, 0.4, "", sz=14, clr=WHITE)
para(tf, "Type:  BEFORE INSERT trigger on REGISTRATION", sz=14, clr=CYAN)
para(tf, "", sz=8)
para(tf, "What it does:", sz=15, clr=WHITE, bold=True, spacing=Pt(8))
para(tf, "▸  Fires automatically before any new registration", sz=14, clr=WHITE)
para(tf, "▸  Counts current registrations for that event", sz=14, clr=WHITE)
para(tf, "▸  Compares count with Max_Seats of the event", sz=14, clr=WHITE)
para(tf, "▸  If event is FULL → ABORTs the insert automatically", sz=14, clr=ORANGE)
para(tf, "", sz=8)
para(tf, "Why it matters:", sz=15, clr=WHITE, bold=True, spacing=Pt(8))
para(tf, "▸  Database-level enforcement, not just application code", sz=14, clr=WHITE)
para(tf, "▸  Impossible to overbook even via direct SQL access", sz=14, clr=WHITE)
para(tf, "▸  Ensures data consistency at all times", sz=14, clr=GREEN)

card(s, 6.8, 1.2, 6.0, 5.5)
txt(s, 7.1, 1.3, 5.4, 0.5, "👁️  View: Participation Summary", sz=20, clr=ACCENT, bold=True)
tf = txt(s, 7.1, 1.9, 5.4, 0.4, "", sz=14, clr=WHITE)
para(tf, "View Name:  vw_event_participation", sz=14, clr=CYAN)
para(tf, "", sz=8)
para(tf, "What it does:", sz=15, clr=WHITE, bold=True, spacing=Pt(8))
para(tf, "▸  A virtual table (no data stored separately)", sz=14, clr=WHITE)
para(tf, "▸  JOINs 3 tables: EVENT + FACULTY + REGISTRATION", sz=14, clr=WHITE)
para(tf, "▸  Uses GROUP BY to count participants per event", sz=14, clr=WHITE)
para(tf, "▸  Shows event name, date, venue, admin, count", sz=14, clr=WHITE)
para(tf, "", sz=8)
para(tf, "Why it matters:", sz=15, clr=WHITE, bold=True, spacing=Pt(8))
para(tf, "▸  Simplifies complex queries into a single access", sz=14, clr=WHITE)
para(tf, "▸  Always returns real-time data from base tables", sz=14, clr=WHITE)
para(tf, "▸  Used in admin dashboard for live statistics", sz=14, clr=GREEN)

# ==================== SLIDE 8: TRANSACTIONS & ACID ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.3, 13.333, 0.6, "🔄  Transactions & ACID Properties", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)

card(s, 0.6, 1.2, 6.0, 3.0)
txt(s, 0.9, 1.3, 5.4, 0.5, "Transaction: Event Registration Flow", sz=18, clr=GOLD, bold=True)
tf = txt(s, 0.9, 1.9, 5.4, 2.2, "", sz=14, clr=WHITE)
para(tf, "Step 1 →  BEGIN transaction", sz=14, clr=CYAN)
para(tf, "Step 2 →  Verify event exists in database", sz=14, clr=WHITE, spacing=Pt(6))
para(tf, "Step 3 →  Insert registration (trigger checks capacity)", sz=14, clr=WHITE, spacing=Pt(6))
para(tf, "Step 4a →  COMMIT if everything succeeds ✓", sz=14, clr=GREEN, spacing=Pt(6))
para(tf, "Step 4b →  ROLLBACK if any error occurs ✗", sz=14, clr=RED, spacing=Pt(6))
para(tf, "", sz=6)
para(tf, "→  Either ALL steps succeed or NONE take effect", sz=14, clr=GOLD, bold=True, spacing=Pt(8))

acid = [
    ("A", "Atomicity", "All or nothing — partial\noperations are impossible"),
    ("C", "Consistency", "Database moves from one\nvalid state to another"),
    ("I", "Isolation", "Concurrent transactions\ndon't interfere with each other"),
    ("D", "Durability", "Once committed, data\nsurvives system crashes"),
]
for i, (letter, name, desc) in enumerate(acid):
    x = 6.8 + i * 1.6
    card(s, x, 1.2, 1.45, 3.0)
    txt(s, x, 1.3, 1.45, 0.6, letter, sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, 1.9, 1.45, 0.35, name, sz=12, clr=GOLD, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, 2.4, 1.25, 1.5, desc, sz=11, clr=MUTED, align=PP_ALIGN.CENTER)

card(s, 0.6, 4.6, 12.2, 2.5)
txt(s, 0.9, 4.7, 11.6, 0.4, "📑  Indexes for Performance Optimization", sz=18, clr=CYAN, bold=True)
tf = txt(s, 0.9, 5.2, 5.5, 1.8, "", sz=14, clr=WHITE)
para(tf, "▸  Index on REGISTRATION(Student_ID)", sz=14, clr=WHITE)
para(tf, "    → Fast lookup of student's registrations", sz=12, clr=MUTED)
para(tf, "▸  Index on REGISTRATION(Event_ID)", sz=14, clr=WHITE, spacing=Pt(6))
para(tf, "    → Fast participant count per event", sz=12, clr=MUTED)
tf2 = txt(s, 6.5, 5.2, 5.5, 1.8, "", sz=14, clr=WHITE)
para(tf2, "▸  Index on EVENT(Faculty_ID)", sz=14, clr=WHITE)
para(tf2, "    → Fast lookup of admin's events", sz=12, clr=MUTED)
para(tf2, "▸  Index on EVENT(Date)", sz=14, clr=WHITE, spacing=Pt(6))
para(tf2, "    → Fast sorting and filtering by date", sz=12, clr=MUTED)

# ==================== SLIDE 9: SQL OPERATIONS USED ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 0.3, 13.333, 0.6, "📝  SQL Operations & Methodology", sz=32, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)

ops_left = [
    ("INSERT", "Adding new students, admins, events, registrations", CYAN),
    ("SELECT", "Querying data with WHERE, LIKE for search/filter", CYAN),
    ("UPDATE", "Modifying event details (name, date, venue, seats)", CYAN),
    ("DELETE", "Removing events, students, and registrations", CYAN),
    ("JOIN", "LEFT JOIN across 3 tables (Event ↔ Faculty ↔ Registration)", ACCENT),
]
ops_right = [
    ("GROUP BY", "Counting participants per event (aggregation)", ACCENT),
    ("SUBQUERY", "Inline participant count within SELECT statements", ACCENT),
    ("TRIGGER", "Automatic overbooking prevention (BEFORE INSERT)", ORANGE),
    ("VIEW", "Virtual table for participation summary", ORANGE),
    ("TRANSACTION", "ACID-compliant BEGIN / COMMIT / ROLLBACK", GREEN),
]

card(s, 0.6, 1.1, 6.0, 5.5)
txt(s, 0.9, 1.2, 5.4, 0.4, "Data Manipulation & Queries", sz=17, clr=GOLD, bold=True)
for j, (op, desc, clr) in enumerate(ops_left):
    y = 1.7 + j * 1.0
    tf = txt(s, 1.0, y, 5.2, 0.9, "", sz=14)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = op; r1.font.color.rgb = clr; r1.font.bold = True; r1.font.size = Pt(16)
    para(tf, desc, sz=13, clr=MUTED, spacing=Pt(2))

card(s, 6.8, 1.1, 6.0, 5.5)
txt(s, 7.1, 1.2, 5.4, 0.4, "Advanced Features", sz=17, clr=GOLD, bold=True)
for j, (op, desc, clr) in enumerate(ops_right):
    y = 1.7 + j * 1.0
    tf = txt(s, 7.2, y, 5.2, 0.9, "", sz=14)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = op; r1.font.color.rgb = clr; r1.font.bold = True; r1.font.size = Pt(16)
    para(tf, desc, sz=13, clr=MUTED, spacing=Pt(2))

# ==================== SLIDE 10: THANK YOU ====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
txt(s, 0, 1.2, 13.333, 0.8, "🎉", sz=56, align=PP_ALIGN.CENTER)
txt(s, 0, 2.2, 13.333, 0.7, "Thank You!", sz=44, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0, 3.2, 13.333, 0.5, "Smart College Event Management System", sz=20, clr=ACCENT, align=PP_ALIGN.CENTER)

stats = [("4", "Tables"), ("BCNF", "Normalized"), ("1", "Trigger"), ("1", "View"), ("4", "Indexes"), ("ACID", "Compliant")]
for i, (val, label) in enumerate(stats):
    x = 1.5 + i * 1.85
    card(s, x, 4.3, 1.6, 1.5)
    txt(s, x, 4.4, 1.6, 0.6, val, sz=24, clr=CYAN, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, 5.1, 1.6, 0.4, label, sz=13, clr=MUTED, align=PP_ALIGN.CENTER)

txt(s, 0, 6.3, 13.333, 0.5, "Questions?", sz=20, clr=MUTED, align=PP_ALIGN.CENTER)

# Save
out = os.path.join(os.path.dirname(__file__), 'DBMS_Project_Presentation.pptx')
prs.save(out)
print(f"✅ Saved: {out}")
